import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from PyQt5.QtWidgets import (QWidget, QToolTip, QPushButton, QApplication, QMessageBox, QDesktopWidget,
                             QMainWindow, QAction, QMenu, QTextEdit, QFileDialog, QHBoxLayout, QLabel,
                             QLineEdit, QGridLayout, QLCDNumber, QSlider, QVBoxLayout, QInputDialog, QSizePolicy, QDialog)
from PyQt5.QtGui import QFont, QIcon, QColor, QPixmap, QScreen
from PyQt5.QtCore import Qt, pyqtSignal, QObject
from pptx.enum.shapes import MSO_CONNECTOR
import sys
from PyQt5 import QtWidgets, QtCore, QtGui
import os
import re
import presentation_builder
import create_dict


class MetroMap(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        QToolTip.setFont(QFont('SansSerif', 8))
        create_btn = QPushButton('Create map', self)
        create_btn.resize(create_btn.sizeHint())
        create_btn.move(50, 50)
        modify_btn = QPushButton('Modify map', self)
        modify_btn.setToolTip("Modify current map")
        modify_btn.resize(modify_btn.sizeHint())
        modify_btn.move(50, 120)
        reset_btn = QPushButton('Reset map', self)
        reset_btn.resize(reset_btn.sizeHint())
        reset_btn.move(50, 190)
        reset_btn.setToolTip("Remove loaded images from map")
        quit_btn = QPushButton('Quit', self)
        quit_btn.move(50, 260)

        vbox = QVBoxLayout()
        vbox.addStretch(1)
        vbox.addWidget(create_btn)
        vbox.addWidget(modify_btn)
        vbox.addWidget(reset_btn)

        hbox = QHBoxLayout()
        hbox.addStretch(1)
        hbox.addLayout(vbox)

        self.setLayout(vbox)
        self.setGeometry(300, 300, 650, 550)
        self.center()
        self.setWindowTitle('Metro Maps')

        create_btn.clicked.connect(self.create_button_clicked)
        modify_btn.clicked.connect(self.modify_button_clicked)
        reset_btn.clicked.connect(self.reset_button_clicked)
        quit_btn.clicked.connect(QApplication.instance().quit)

    def center(self):
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())

    def create_button_clicked(self):
        if not hasattr(self, 'CreateMap'):
            self.window = CreateMap(self)
        self.window.show()
        self.hide()

    def modify_button_clicked(self):
        presentation_builder.build()
        build_map()
        if not hasattr(self, 'PowerPoint'):
            self.pp = PowerPoint(self)
        self.pp.show()
        self.hide()

    def reset_button_clicked(self):
        create_dict.create_json_dicts()
        presentation_builder.build()


class CreateMap(QMainWindow):
    def __init__(self, parent):
        super(CreateMap, self).__init__(parent)
        self.setAttribute(QtCore.Qt.WA_DeleteOnClose)
        self.initUI()

    def initUI(self):
        pic_btn = QPushButton("Add picture", self)
        pic_btn.move(40, 50)
        pic_btn.resize(pic_btn.sizeHint())

        text_btn = QPushButton("Add text", self)
        text_btn.move(40, 120)
        text_btn.resize(text_btn.sizeHint())

        done_btn = QPushButton("Back", self)
        done_btn.move(40, 190)
        done_btn.resize(done_btn.sizeHint())

        pic_btn.clicked.connect(self.pic_btn_clicked)
        text_btn.clicked.connect(self.text_btn_clicked)
        done_btn.clicked.connect(self.done_btn_clicked)

        self.statusBar()
        self.setGeometry(300, 300, 650, 550)
        self.setWindowTitle('Create new map')

        vbox = QVBoxLayout()
        vbox.addStretch(1)
        vbox.addWidget(pic_btn)
        vbox.addWidget(text_btn)
        vbox.addWidget(done_btn)

        hbox = QHBoxLayout()
        hbox.addStretch(1)
        hbox.addLayout(vbox)

        self.setLayout(vbox)
        self.center()

    def pic_btn_clicked(self):
        self.showMinimized()
        os.system("START Snipping_Tool.lnk")

    def text_btn_clicked(self):
        if not hasattr(self, 'TextEdit'):
            self.text = TextEdit(self)
        self.text.show()
        self.hide()

    def done_btn_clicked(self):
        self.close()
        self.parent().show()

    def center(self):
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())


class TextEdit(QMainWindow):
    def __init__(self, parent):
        super(TextEdit, self).__init__(parent)
        self.setAttribute(QtCore.Qt.WA_DeleteOnClose)
        self.initUI()
        self.num_shot = 0

    def initUI(self):
        try:
            self.textEdit = QTextEdit()
            self.setCentralWidget(self.textEdit)
            save_pic = QAction('Save and Exit', self)
            save_pic.setShortcut('Ctrl+S')
            save_pic.triggered.connect(self.take_screenshot)
            go_back = QAction('Back', self)
            go_back.triggered.connect(self.go_back)

            menu_bar = self.menuBar()
            file_menu = menu_bar.addMenu('&File')
            file_menu.addAction(save_pic)
            file_menu.addAction(go_back)

            self.setGeometry(300, 300, 650, 600)
            self.setWindowTitle('Text Editor')
            self.center()
            self.show()
        except Exception as e:
            print(e)

    def take_screenshot(self):
        try:
            self.num_shot += 1
            p = app.primaryScreen().grabWindow(self.winId())
            if not os.path.isdir(os.path.join(os.path.dirname(os.path.realpath(__file__)), "{}\{}".format("dist", "text"))):
                os.makedirs(os.path.join(os.path.dirname(os.path.realpath(__file__)), "{}\{}".format("dist", "text")))
            else:
                file_list = os.listdir(os.path.join(os.path.dirname(os.path.realpath(__file__)), "{}\{}".format("dist", "text")))
                if file_list:
                    regex = re.compile(r'\d+')
                    num_list = []
                    for filename in file_list:
                        num_list.append(int(regex.findall(filename)[0]))
                    self.num_shot = max(num_list) + 1
            filename = "text_{}".format(self.num_shot)
            root = os.path.join(os.path.dirname(os.path.realpath(__file__)), "{}\{}".format("dist", "text"))
            p.save("{}.{}".format(os.path.join(root, filename), 'jpg'))
            self.hide()
            self.parent().show()
        except Exception as e:
            print(e)

    def center(self):
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())

    def go_back(self):
        self.hide()
        self.parent().show()


class MessageBox(QDialog):
    def __init__(self, parent=None):
        super(MessageBox, self).__init__(parent)
        msgBox = QMessageBox()
        msgBox.setText('Which line?')
        msgBox.addButton(QPushButton('Red'), QMessageBox.YesRole)
        msgBox.addButton(QPushButton('Green'), QMessageBox.NoRole)
        msgBox.addButton(QPushButton('Blue'), QMessageBox.RejectRole)
        self.ret = msgBox.exec_()


class NodeMessageBox(QDialog):
    def __init__(self, parent=None):
        super(NodeMessageBox, self).__init__(parent)
        msgBox = QMessageBox()
        msgBox.setText('Node?')
        msgBox.addButton(QPushButton('1'), QMessageBox.YesRole)
        msgBox.addButton(QPushButton('2'), QMessageBox.AcceptRole)
        msgBox.addButton(QPushButton('3'), QMessageBox.NoRole)
        msgBox.addButton(QPushButton('4'), QMessageBox.RejectRole)
        self.ret = msgBox.exec_()


class DestLineMessageBox(QDialog):
    def __init__(self, parent=None):
        super(DestLineMessageBox, self).__init__(parent)
        msgBox = QMessageBox()
        msgBox.setText('To which line?')
        msgBox.addButton(QPushButton('Red'), QMessageBox.YesRole)
        msgBox.addButton(QPushButton('Green'), QMessageBox.NoRole)
        msgBox.addButton(QPushButton('Blue'), QMessageBox.RejectRole)
        self.ret = msgBox.exec_()


class PowerPoint(QMainWindow):
    def __init__(self, parent):
        super(PowerPoint, self).__init__(parent)
        self.initUI()
        self.setAttribute(QtCore.Qt.WA_DeleteOnClose)

    def initUI(self):
        openFile = QPushButton('Open file', self)
        openFile.setToolTip('Choose a file to put on map')
        openFile.clicked.connect(self.showDialog)
        openFile.resize(openFile.sizeHint())
        openFile.move(50, 50)
        draw_lines = QPushButton('Draw lines', self)
        openFile.setToolTip('Draw lines between nodes')
        draw_lines.clicked.connect(self.draw_lines_clicked)
        draw_lines.resize(draw_lines.sizeHint())
        draw_lines.move(50, 120)
        open_prs = QPushButton('Load map', self)
        open_prs.resize(open_prs.sizeHint())
        open_prs.move(50, 190)
        open_prs.clicked.connect(open_ppt)
        back_btn = QPushButton('Back', self)
        back_btn.resize(back_btn.sizeHint())
        back_btn.move(50, 260)
        back_btn.clicked.connect(self.back_btn_clicked)

        self.setGeometry(300, 300, 650, 550)
        self.setWindowTitle('Modify')
        self.center()

    def back_btn_clicked(self):
        self.close()
        self.parent().show()

    def draw_lines_clicked(self):
        try:
            with open("data.json", 'r') as f:
                mapping = json.load(f)
            color_dict = {0: "Red", 1: "Green", 2: "Blue"}
            prs = Presentation('output.pptx')
            slide = prs.slides[0]
            lineReply = MessageBox()
            lineReply = color_dict[lineReply.ret]
            nodeReply = NodeMessageBox()
            destLineReply = DestLineMessageBox()
            destLineReply = color_dict[destLineReply.ret]
            destNodeReply = NodeMessageBox()

            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mapping['line_mapping'][lineReply][nodeReply.ret][0],
                                              mapping['line_mapping'][lineReply][nodeReply.ret][1],
                                              mapping['line_mapping'][destLineReply][destNodeReply.ret][0],
                                              mapping['line_mapping'][destLineReply][destNodeReply.ret][1])
            mapping['lines'].append([mapping['line_mapping'][lineReply][nodeReply.ret][0],
                                     mapping['line_mapping'][lineReply][nodeReply.ret][1],
                                     mapping['line_mapping'][destLineReply][destNodeReply.ret][0],
                                     mapping['line_mapping'][destLineReply][destNodeReply.ret][1], lineReply])
            presentation_builder.get_line(lineReply, line, True)
            slide.shapes._spTree.remove(line._element)
            slide.shapes._spTree.insert(2, line._element)
            prs.save("output.pptx")
            with open("data.json", 'w') as f:
                json.dump(mapping, f)
            self.parent().show()
        except Exception as e:
            print(e)

    def showDialog(self):
        with open("data.json", 'r') as f:
            mapping = json.load(f)
        prs = Presentation('output.pptx')
        slide = prs.slides[0]
        fname = QFileDialog.getOpenFileName(self, 'Open file', '/home')
        if fname[0]:
            buttonReply = MessageBox()
            if buttonReply.ret == 0:
                clicked = 'Red'
            elif buttonReply.ret == 1:
                clicked = 'Green'
            else:
                clicked = 'Blue'
            for i in range(0, 4):
                if not mapping['pic_mapping'][clicked][i][4]:
                    mapping['pic_mapping'][clicked][i][4] = True
                    mapping['pic_list'][clicked][i] = fname[0]
                    pic = slide.shapes.add_picture(fname[0], mapping['pic_mapping'][clicked][i][0],
                                                   mapping['pic_mapping'][clicked][i][1],
                                                   mapping['pic_mapping'][clicked][i][2],
                                                   mapping['pic_mapping'][clicked][i][3])
                    pic.click_action.hyperlink.address = fname[0]
                    with open("data.json", 'w') as f:
                        json.dump(mapping, f)
                    break
        prs.save('output.pptx')
        self.parent().show()

    def center(self):
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())


def open_ppt():
    os.system('output.pptx')


def build_map():
    try:
        with open("data.json", 'r') as f:
            mapping = json.load(f)
        prs = Presentation('output.pptx')
        slide = prs.slides[0]
        for i in ["Red", "Green", "Blue"]:
            for j in range(0, 4):
                if mapping['pic_list'][i][j]:
                    pic = slide.shapes.add_picture(mapping['pic_list'][i][j], mapping['pic_mapping'][i][j][0],
                                                   mapping['pic_mapping'][i][j][1],
                                                   mapping['pic_mapping'][i][j][2],
                                                   mapping['pic_mapping'][i][j][3])
                    pic.click_action.hyperlink.address = mapping['pic_list'][i][j]
        if mapping['lines']:
            for connection in mapping['lines']:
                line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, connection[0], connection[1], connection[2],
                                                  connection[3])
                presentation_builder.get_line(connection[4], line, True)
                slide.shapes._spTree.remove(line._element)
                slide.shapes._spTree.insert(2, line._element)
        prs.save("output.pptx")
        return
    except Exception as e:
        print(e)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    map = MetroMap()
    map.show()
    app.aboutToQuit.connect(app.deleteLater)
    sys.exit(app.exec_())

