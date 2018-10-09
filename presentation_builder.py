import os
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.action import PP_ACTION
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector


def get_line(color, connector, bright=None):
    try:
        colors_dict = {"Red": RGBColor(255, 0, 0), "Green": RGBColor(0, 255, 0), "Blue": RGBColor(0, 0, 255)}
        line = LineFormat(connector)
        line.fill.solid()
        line.fill.fore_color.rgb = colors_dict["{}".format(color)]
        if bright:
            line.color.brightness = 0.7
        else:
            line.color.brightness = 0.5
        line.width = Pt(7)
        return line.fill.fore_color.rgb
    except Exception as e:
        print(e)


def draw_nodes(slide):
    try:
        width, height = Inches(0.2), Inches(0.2)
        top = Inches(1.4)
        for y in range(0, 3):
            left = 0
            for x in range(0, 4):
                left += Inches(2)
                # top += Inches(1.4)
                x = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
                fill = x.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
            left = 0
            top += Inches(2)
        return
    except Exception as e:
        print(e)


def build():
    LAYOUT_INDEX = 6
    root = os.path.dirname(os.path.realpath(__file__))

    # create presentation with 1 slide ------
    begin_x = Inches(1)
    begin_y = Inches(1.5)
    end_x = Inches(9)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_INDEX])

    # Draw 3 lines.
    connector_first = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, begin_y)
    connector_second = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y + Inches(2), end_x, begin_y + Inches(2))
    connector_third = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y + Inches(4), end_x, begin_y + Inches(4))

    get_line("Red", connector_first)
    get_line("Green", connector_second)
    get_line("Blue", connector_third)
    draw_nodes(slide)

    prs.save('./output.pptx')


# build()
# prs.save(os.path.join(root, 'test.pptx'))


